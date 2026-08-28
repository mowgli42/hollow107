#!/usr/bin/env python3
"""Generate Brain Book whitepaper, sponsorship memo, and phase deck."""

from pathlib import Path

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_LINE_SPACING
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor, Twips
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import nsmap
from pptx.util import Emu, Inches as PptInches, Pt as PptPt

INK = RGBColor(0x1A, 0x1C, 0x1F)
STEEL = RGBColor(0x3D, 0x4F, 0x5F)
MUTED = RGBColor(0x5C, 0x5F, 0x66)
RULE = RGBColor(0xC9, 0xC2, 0xB2)
PAPER = RGBColor(0xF3, 0xEF, 0xE6)

P_INK = PptRGB(0x1A, 0x1C, 0x1F)
P_STEEL = PptRGB(0x3D, 0x4F, 0x5F)
P_MUTED = PptRGB(0x5C, 0x5F, 0x66)
P_PAPER = PptRGB(0xF3, 0xEF, 0xE6)
P_ELEV = PptRGB(0xFA, 0xF8, 0xF2)
P_WHITE = PptRGB(0xFF, 0xFF, 0xFF)
P_RULE = PptRGB(0xC9, 0xC2, 0xB2)

OUT_PUBLIC = Path("/workspace/public/docs")
OUT_ART = Path("/workspace/artifacts")


def set_run_font(run, name="Calibri", size=11, bold=False, italic=False, color=INK):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    r = run._element
    rPr = r.get_or_add_rPr()
    rFonts = rPr.find(qn("w:rFonts"))
    if rFonts is None:
        rFonts = OxmlElement("w:rFonts")
        rPr.append(rFonts)
    rFonts.set(qn("w:ascii"), name)
    rFonts.set(qn("w:hAnsi"), name)
    rFonts.set(qn("w:eastAsia"), name)
    rFonts.set(qn("w:cs"), name)


def shade_cell(cell, hex_color):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), hex_color)
    shd.set(qn("w:val"), "clear")
    tcPr.append(shd)


def set_cell_border(cell, color="C9C2B2"):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    tcBorders = OxmlElement("w:tcBorders")
    for edge in ("top", "left", "bottom", "right"):
        el = OxmlElement(f"w:{edge}")
        el.set(qn("w:val"), "single")
        el.set(qn("w:sz"), "4")
        el.set(qn("w:space"), "0")
        el.set(qn("w:color"), color)
        tcBorders.append(el)
    tcPr.append(tcBorders)


def add_page_number(paragraph):
    run1 = paragraph.add_run("INTERNAL CONCEPT BRIEF  ·  ")
    set_run_font(run1, size=9, color=MUTED)
    run2 = paragraph.add_run("Page ")
    set_run_font(run2, size=9, color=MUTED)
    fldChar1 = OxmlElement("w:fldChar")
    fldChar1.set(qn("w:fldCharType"), "begin")
    instr = OxmlElement("w:instrText")
    instr.set(qn("xml:space"), "preserve")
    instr.text = " PAGE "
    fldChar2 = OxmlElement("w:fldChar")
    fldChar2.set(qn("w:fldCharType"), "end")
    r = paragraph.add_run()
    set_run_font(r, size=9, color=MUTED)
    r._r.append(fldChar1)
    r2 = paragraph.add_run()
    set_run_font(r2, size=9, color=MUTED)
    r2._r.append(instr)
    r3 = paragraph.add_run()
    set_run_font(r3, size=9, color=MUTED)
    r3._r.append(fldChar2)


def add_body(doc, text, size=11, space_after=10, italic=False):
    p = doc.add_paragraph()
    p.paragraph_format.space_after = Pt(space_after)
    p.paragraph_format.space_before = Pt(0)
    p.paragraph_format.line_spacing = 1.15
    run = p.add_run(text)
    set_run_font(run, size=size, italic=italic)
    return p


def add_heading_custom(doc, text, size=18):
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(18)
    p.paragraph_format.space_after = Pt(8)
    run = p.add_run(text)
    set_run_font(run, name="Georgia", size=size, bold=True, color=INK)
    return p


def add_kicker(doc, text):
    p = doc.add_paragraph()
    p.paragraph_format.space_before = Pt(16)
    p.paragraph_format.space_after = Pt(0)
    run = p.add_run(text.upper())
    set_run_font(run, size=9, bold=True, color=STEEL)
    return p


def bullet(doc, text, bold_lead=None):
    p = doc.add_paragraph(style="List Bullet")
    p.paragraph_format.space_after = Pt(4)
    p.paragraph_format.left_indent = Inches(0.35)
    if bold_lead:
        r = p.add_run(bold_lead)
        set_run_font(r, size=11, bold=True)
        r2 = p.add_run(text)
        set_run_font(r2, size=11)
    else:
        r = p.add_run(text)
        set_run_font(r, size=11)
    return p


def build_whitepaper(path: Path):
    doc = Document()
    section = doc.sections[0]
    section.page_width = Inches(8.5)
    section.page_height = Inches(11)
    section.left_margin = Inches(1.1)
    section.right_margin = Inches(1.1)
    section.top_margin = Inches(1.0)
    section.bottom_margin = Inches(1.0)

    header = section.header
    header.is_linked_to_previous = False
    hp = header.paragraphs[0]
    hp.alignment = WD_ALIGN_PARAGRAPH.LEFT
    r = hp.add_run("BRAIN BOOK  ·  Internal concept brief  ·  August 2026")
    set_run_font(r, size=9, color=STEEL)

    footer = section.footer
    footer.is_linked_to_previous = False
    fp = footer.paragraphs[0]
    fp.alignment = WD_ALIGN_PARAGRAPH.LEFT
    add_page_number(fp)
    r = fp.add_run("  ·  Not a technical order  ·  Not an airworthy disposition")
    set_run_font(r, size=9, color=MUTED)

    k = doc.add_paragraph()
    k.paragraph_format.space_after = Pt(6)
    run = k.add_run("WHITEPAPER")
    set_run_font(run, size=10, bold=True, color=STEEL)

    title = doc.add_paragraph()
    title.paragraph_format.space_after = Pt(10)
    run = title.add_run(
        "The Brain Book: capturing integration knowledge before Field Service has to reinvent it"
    )
    set_run_font(run, name="Georgia", size=22, bold=True)

    add_body(
        doc,
        "A local, cited knowledge compiler for Field Service Representatives and, later, Air Force T.O. 00-25-107 technical assistance requests. Procedures remain airworthy. The Brain Book is the evidence and the memory those procedures cannot hold.",
        size=12,
        italic=True,
        space_after=14,
    )

    add_kicker(doc, "01  ·  The need")
    add_heading_custom(doc, "Knowledge is created in build and lost before the field", 16)
    add_body(
        doc,
        "A new aircraft does not wait for a technical order to accumulate knowledge. Supplier reports, lab integration, software loads, tickets, and flight-test cards already describe what the system does when it is healthy and what it looks like when it is not. That body of knowledge is real. It is also perishable. It lives in people, in ticket threads, and in log files that nobody indexes.",
    )
    add_body(
        doc,
        "The only official path to the field today is a written procedure. Procedures are necessary. They are also late, flattened, and config-blind. They cannot carry the log signature of a cold-soak align failure, the fact that a given serial number had an inverted timing pulse until an ECO, or the difference between two OFP loads on the same symptom. When a Field Service Representative (FSR) later faces that symptom — and when the unit files a 107 — they start from the procedure and from whoever still remembers last year’s rack.",
    )
    add_body(
        doc,
        "The system already produces the missing evidence: logs, BIT, and event traces. We do not parse them against known-good baselines. We do not keep the pairing of this ticket, this bad log, this good log, this software hash. That pairing is the Brain Book. Without it, every inbound 107 becomes a manual reconstruction of work the program already paid for.",
    )

    add_kicker(doc, "02  ·  What it is")
    add_heading_custom(doc, "A versioned book of reviewed objects, not a chatbot", 16)
    add_body(
        doc,
        "A Brain Book is not a chatbot over manuals and it is not a replacement for a TO. It is a versioned folder of reviewed objects, tied to MDS, block, OFP, and ICD:",
    )
    bullet(
        doc,
        " — log strings and BIT codes mapped to source file and meaning.",
        "Event dictionary",
    )
    bullet(
        doc,
        " — what a healthy power-up, align, or test card emits on a given OFP.",
        "Mode baselines",
    )
    bullet(
        doc,
        " — IS / IS-NOT, events only in the bad log, distinguishing check, resolution, citations.",
        "Signatures",
    )
    bullet(
        doc,
        " — the narrative labs never put in a procedure: what was tried, what failed, what closed it.",
        "Ticket digests",
    )
    add_body(
        doc,
        "Drafts are generated by a local LLM from artifacts the program already has: source, associated documents, tickets, and paired good/bad logs. Nothing is treated as true until a lab, flight-test, or software owner changes maturity from draft to reviewed. Citations are mandatory. Ungrounded numbers are refused. The official 107 response and any airworthy step remain human.",
    )

    add_kicker(doc, "03  ·  Consumers")
    add_heading_custom(doc, "FSRs first; 107s when the process stands up", 16)
    add_body(
        doc,
        "We are not yet in 107 land. That is the point. FSRs and integration engineers are the first consumers. They need a way to ask “have we seen this?” against the same configuration that is on the jet, and to turn a log dump into a short list of matching signatures instead of a scroll through raw files.",
    )
    add_body(
        doc,
        "When 107s (T.O. 00-25-107 Technical Assistance Requests and Maintenance Assistance Requests) begin to arrive, the same book becomes the triage layer. Most inbound requests fail first on incomplete facts: missing OFP, no last-known-good, no “already tried,” no log attached. The model’s first job is not diagnosis. It is to name the missing context, draft the callback questions, and only then rank hypotheses that the signatures can support.",
    )
    add_body(
        doc,
        "That sequence — generate memory now, use it for FSR troubleshooting next, feed 107 triage later — is how knowledge stops walking off the aircraft with the people who built it.",
    )

    add_kicker(doc, "04  ·  Timing")
    add_heading_custom(doc, "The compiler has to run while the logs are still attached", 16)
    add_body(
        doc,
        "A language model cannot recover a signature that was never captured. If a ticket closes with “fixed it, see verbal,” the field has nothing to match. Generation must happen while supplier reports, lab anomalies, and flight-test cards are still attached to logs.",
    )
    add_body(
        doc,
        "The near-term product is therefore a Brain Book factory, not a field portal. One subsystem, one OFP, historical tickets with logs, and a human review loop is enough to prove the idea. If engineers who chased those problems recognize the drafted signatures, the later 107 workflow is just a new cover sheet on the same folder.",
    )

    add_kicker(doc, "05  ·  Three phases")
    add_heading_custom(doc, "One book, three operating pictures", 16)

    table = doc.add_table(rows=4, cols=3)
    table.autofit = True
    headers = ["Phase", "What happens", "Output"]
    rows = [
        [
            "1. Generate",
            "Ingest source, documents, tickets, and paired good/bad logs. Deterministic decode and diff; local LLM drafts signatures and digests.",
            "Reviewed dictionary, baselines, signatures (maturity: draft → lab/FT reviewed).",
        ],
        [
            "2. 107 triage",
            "Normalize the request. Completeness gate on configuration and evidence. Retrieve only for the stated OFP/block. Block diagnosis if facts are missing.",
            "Callback questions, or a triage packet with top-three hypotheses and kill-checks.",
        ],
        [
            "3. FSR log review",
            "Decode logs to events. Match signatures. LLM interprets only unmatched residue. One next check. Update IS/IS-NOT and fishbone status.",
            "Case file: matches, unmatched events, next check, citations. New signature if the case is novel.",
        ],
    ]
    for i, h in enumerate(headers):
        cell = table.rows[0].cells[i]
        cell.text = ""
        p = cell.paragraphs[0]
        run = p.add_run(h)
        set_run_font(run, size=10, bold=True, color=RGBColor(0xF3, 0xEF, 0xE6))
        shade_cell(cell, "1A1C1F")
        set_cell_border(cell, "1A1C1F")
    for r_i, row in enumerate(rows):
        for c_i, val in enumerate(row):
            cell = table.rows[r_i + 1].cells[c_i]
            cell.text = ""
            p = cell.paragraphs[0]
            run = p.add_run(val)
            set_run_font(run, size=10, bold=(c_i == 0))
            shade_cell(cell, "FAF8F2" if r_i % 2 == 0 else "F3EFE6")
            set_cell_border(cell)

    add_body(doc, "", space_after=6)
    add_body(
        doc,
        "Log parsing is software first. Do not point the model at raw dumps. Decode to an event table, match against the signature library, and use the LLM only on the residue and on incomplete free text.",
    )

    add_kicker(doc, "06  ·  Guardrails")
    add_heading_custom(doc, "Advise. Do not dispose. Do not command.", 16)
    add_body(
        doc,
        "Local inference only — source, tickets, and logs do not leave the enclave. Decode and signature match are deterministic; the model writes names, IS / IS-NOT, and open questions from provided text. Live commanding of hardware is out of scope. The Brain Book advises. The signed procedure and the FSR dispose.",
    )

    add_kicker(doc, "07  ·  Recommended next step")
    add_heading_custom(doc, "A bounded demo, not a platform", 16)
    add_body(
        doc,
        "Sponsor a small part-time team to compile a Brain Book for one subsystem and one OFP using an open local stack (NVIDIA Nemotron via Ollama, Open WebUI, Python/ripgrep parsers, Git-backed Markdown/YAML). Success is ten reviewed signatures that the engineers who chased those problems recognize, plus historical cases where the eventual cause appears in the top three after decode. Detail is in the companion sponsorship memo.",
    )

    path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(path))


def build_memo(path: Path):
    doc = Document()
    section = doc.sections[0]
    section.page_width = Inches(8.5)
    section.page_height = Inches(11)
    section.left_margin = Inches(1.1)
    section.right_margin = Inches(1.1)
    section.top_margin = Inches(1.0)
    section.bottom_margin = Inches(0.9)

    header = section.header
    hp = header.paragraphs[0]
    r = hp.add_run("BRAIN BOOK  ·  Sponsorship memorandum  ·  August 2026")
    set_run_font(r, size=9, color=STEEL)

    footer = section.footer
    fp = footer.paragraphs[0]
    r = fp.add_run("Internal  ·  Concept only  ·  Not a technical order")
    set_run_font(r, size=9, color=MUTED)

    t = doc.add_paragraph()
    run = t.add_run("Memorandum")
    set_run_font(run, name="Georgia", size=22, bold=True)

    meta = [
        ("TO:", "Program leadership / funding sponsor"),
        ("FROM:", "Integration / Field Service concept team"),
        ("DATE:", "28 August 2026"),
        ("SUBJECT:", "Part-time sponsorship to compile a Brain Book (local Nemotron demo)"),
    ]
    for k, v in meta:
        p = doc.add_paragraph()
        p.paragraph_format.space_after = Pt(2)
        r1 = p.add_run(k + "  ")
        set_run_font(r1, size=11, bold=True, color=STEEL)
        r2 = p.add_run(v)
        set_run_font(r2, size=11)

    line = doc.add_paragraph()
    line.paragraph_format.space_before = Pt(8)
    line.paragraph_format.space_after = Pt(10)
    run = line.add_run("—" * 36)
    set_run_font(run, size=8, color=RULE)

    paras = [
        "We are losing integration knowledge as this aircraft is built. Supplier reports, lab bring-up, tickets, and flight-test logs already describe what healthy and failed look like, but the only extract to the field is a procedure written later by hand. Field Service Representatives will be asked to diagnose from that thin slice, and when 107s (T.O. 00-25-107 assistance requests) start arriving they will reconstruct work the program already paid for. We need a Brain Book compiled now — event dictionary, good-log baselines, and fault signatures — before the people and the logs scatter.",
        "We are asking sponsorship for a small part-time team to prove Brain Book generation and FSR-style log triage on one subsystem and one OFP. The demo will ingest source (log/BIT strings), associated documents, tickets, and paired good/bad logs; emit draft signatures and ticket digests; and require a human review stamp before anything is treated as true. When a sample field report plus a log is dropped on that book, the demo will first list missing context, then match signatures and recommend one next check. It will not write official dispositions or command hardware.",
        "The stack stays open and local: NVIDIA Nemotron (Nemotron 3.5 Lightning / Nano class, ~30B with ~3B active) served on an existing GPU workstation via Ollama; Open WebUI for the operator; Python decoders and ripgrep for logs and source; Git-backed Markdown/YAML for the book. No cloud tokens, no new platform, no Wiki or observability suite in v1. Nemotron is open-weight, agent-oriented, and small enough to keep source and logs on the box.",
        "Ask: two to three people part-time for eight to twelve weeks, one workstation, and access to a bounded corpus (one repo, the ICD/ATP pages, ~20 tickets with logs). Success is ten reviewed signatures that the engineers who chased those problems recognize, plus side-by-side triage of historical cases where the eventual cause appears in the top three after decode. If that holds, the same folder becomes the FSR and 107 front end without a second architecture.",
    ]
    for para in paras:
        add_body(doc, para, size=11, space_after=12)

    add_heading_custom(doc, "Proposed stack", 14)
    table = doc.add_table(rows=6, cols=3)
    for i, h in enumerate(["Layer", "Choice", "Why"]):
        cell = table.rows[0].cells[i]
        cell.text = ""
        p = cell.paragraphs[0]
        run = p.add_run(h)
        set_run_font(run, size=10, bold=True, color=RGBColor(0xF3, 0xEF, 0xE6))
        shade_cell(cell, "1A1C1F")
        set_cell_border(cell, "1A1C1F")
    data = [
        ["Model", "Nemotron 3.5 Lightning via Ollama", "Open weights, local, agentic, workstation-sized"],
        ["UI", "Open WebUI", "Chat against the book; no custom app in v1"],
        ["Decode", "Python parsers + ripgrep", "Logs and source stay deterministic"],
        ["Store", "Git Markdown / YAML", "Reviewable, versioned to OFP"],
        ["Out of scope", "Cloud LLM, Wiki, LGTM, MCP mesh", "Add only after signatures are trusted"],
    ]
    for r_i, row in enumerate(data):
        for c_i, val in enumerate(row):
            cell = table.rows[r_i + 1].cells[c_i]
            cell.text = ""
            p = cell.paragraphs[0]
            run = p.add_run(val)
            set_run_font(run, size=10, bold=(c_i == 0))
            shade_cell(cell, "FAF8F2" if r_i % 2 == 0 else "F3EFE6")
            set_cell_border(cell)

    add_body(doc, "", space_after=8)
    add_heading_custom(doc, "Decision requested", 14)
    add_body(
        doc,
        "Approve part-time use of two to three people and one GPU workstation for eight to twelve weeks, with access to a single-subsystem corpus, to produce a reviewed Brain Book slice and a log-triage demonstration against historical cases.",
    )

    path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(path))


# --- PPTX helpers ---

def _set_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def _set_line(shape, rgb, pt=1.0):
    shape.line.color.rgb = rgb
    shape.line.width = PptPt(pt)


def _tf(shape, text, size=14, bold=False, color=P_INK, align=PP_ALIGN.LEFT, font="Calibri"):
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.clear()
    run = p.add_run()
    run.text = text
    run.font.size = PptPt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return tf


def _add_text(slide, l, t, w, h, text, size=14, bold=False, color=P_INK, align=PP_ALIGN.LEFT, font="Calibri", anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(PptInches(l), PptInches(t), PptInches(w), PptInches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf._txBody.bodyPr.set("anchor", {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}.get(anchor, "t"))
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = PptPt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return box


def _card(slide, l, t, w, h, title, body, fill=P_ELEV, title_color=P_INK, body_color=P_MUTED, ink_bar=False):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptInches(l), PptInches(t), PptInches(w), PptInches(h))
    _set_fill(sh, fill)
    _set_line(sh, P_RULE, 0.75)
    try:
        sh.adjustments[0] = 0.08
    except Exception:
        pass
    _add_text(slide, l + 0.18, t + 0.12, w - 0.36, 0.36, title, size=13, bold=True, color=title_color)
    _add_text(slide, l + 0.18, t + 0.46, w - 0.36, h - 0.58, body, size=12, color=body_color)
    return sh


def _footer(slide, page, total=14):
    _add_text(slide, 0.6, 7.12, 8, 0.28, "BRAIN BOOK  ·  INTERNAL  ·  CONCEPT", size=10, color=P_MUTED)
    _add_text(slide, 11.4, 7.12, 1.3, 0.28, f"{page}  /  {total}", size=10, color=P_MUTED, align=PP_ALIGN.RIGHT)


def _kicker_title(slide, kicker, title, subtitle=None):
    _add_text(slide, 0.6, 0.28, 12, 0.28, kicker.upper(), size=11, bold=True, color=P_STEEL)
    _add_text(slide, 0.6, 0.52, 12, 0.7, title, size=28, bold=True, color=P_INK, font="Georgia")
    if subtitle:
        _add_text(slide, 0.6, 1.22, 12, 0.5, subtitle, size=14, color=P_MUTED)


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _bg(slide, color=P_PAPER):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def build_deck(path: Path):
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)
    total = 14

    # 1 Title
    s = _blank(prs)
    _bg(s, P_INK)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(0), PptInches(0), PptInches(0.16), PptInches(7.5))
    _set_fill(bar, P_STEEL)
    _add_text(s, 0.7, 1.6, 11, 0.35, "INTERNAL CONCEPT BRIEF  ·  AUGUST 2026", size=12, bold=True, color=P_STEEL)
    _add_text(s, 0.7, 2.05, 12, 1.6, "The Brain Book", size=48, bold=True, color=P_PAPER, font="Georgia")
    _add_text(
        s,
        0.7,
        3.7,
        11,
        1.0,
        "Capturing lab and flight-test knowledge so Field Service\nand later 107s are not starting from folklore.",
        size=18,
        color=P_PAPER,
    )
    _add_text(s, 0.7, 6.6, 11, 0.35, "Not a technical order  ·  Not an airworthy disposition", size=12, color=P_MUTED)

    # 2 Problem
    s = _blank(prs)
    _bg(s)
    _kicker_title(s, "The leak", "Knowledge is created as the jet is built.\nIt does not travel with the jet.")
    items = [
        ("Created", "Supplier reports, source, lab racks, tickets, flight-test cards, good and bad logs."),
        ("Published", "A procedure, late and flattened. Config-blind. No log signature."),
        ("Lost", "The pairing of ticket + OFP + bad log + good log never leaves the lab."),
    ]
    for i, (t, b) in enumerate(items):
        x = 0.6 + i * 4.15
        sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptInches(x), PptInches(2.6), PptInches(3.9), PptInches(3.5))
        _set_fill(sh, P_INK if i == 2 else P_ELEV)
        _set_line(sh, P_INK if i == 2 else P_RULE, 0.75)
        tc = P_PAPER if i == 2 else P_INK
        bc = PptRGB(0xC9, 0xC2, 0xB2) if i == 2 else P_MUTED
        _add_text(s, x + 0.25, 2.8, 3.4, 0.4, f"0{i+1}", size=12, bold=True, color=P_STEEL if i != 2 else P_STEEL)
        _add_text(s, x + 0.25, 3.2, 3.4, 0.5, t, size=22, bold=True, color=tc, font="Georgia")
        _add_text(s, x + 0.25, 3.85, 3.4, 1.8, b, size=15, color=bc)
    _footer(s, 2, total)

    # 3 Why procedures fail
    s = _blank(prs)
    _bg(s)
    _kicker_title(
        s,
        "Why a TO is not enough",
        "Procedures stay. They cannot hold the evidence.",
        "Airworthy steps still live in the signed book. The Brain Book holds what the TO flattens away.",
    )
    rows = [
        ("Late", "Written after the learning happened."),
        ("Flattened", "One symptom → one action. No OFP split."),
        ("Silent on logs", "Does not say what the BIT stream looked like."),
        ("No IS / IS-NOT", "Does not say which SN, mode, or site it is not."),
    ]
    for i, (t, b) in enumerate(rows):
        y = 2.15 + i * 1.1
        n = s.shapes.add_shape(MSO_SHAPE.OVAL, PptInches(0.7), PptInches(y + 0.12), PptInches(0.42), PptInches(0.42))
        _set_fill(n, P_INK)
        _add_text(s, 0.7, y + 0.16, 0.42, 0.36, str(i + 1), size=12, bold=True, color=P_PAPER, align=PP_ALIGN.CENTER)
        _add_text(s, 1.4, y, 10.5, 0.4, t, size=18, bold=True, color=P_INK, font="Georgia")
        _add_text(s, 1.4, y + 0.4, 10.5, 0.4, b, size=14, color=P_MUTED)
    _footer(s, 3, total)

    # 4 What it is
    s = _blank(prs)
    _bg(s)
    _kicker_title(s, "Definition", "Four objects, versioned to OFP.")
    cards = [
        ("Event dictionary", "Log strings and BIT codes mapped to file:line and meaning. Built from source first."),
        ("Mode baselines", "What healthy power-up, align, or a test card emits on this OFP."),
        ("Signatures", "IS / IS-NOT, events only in the bad log, distinguishing check, resolution, citations."),
        ("Ticket digests", "What was tried, what failed, what closed it — the narrative a TO never keeps."),
    ]
    for i, (t, b) in enumerate(cards):
        x = 0.6 + (i % 2) * 6.3
        y = 2.15 + (i // 2) * 2.2
        _card(s, x, y, 6.0, 2.0, t, b, fill=P_ELEV)
    _footer(s, 4, total)

    # 5 Three phases
    s = _blank(prs)
    _bg(s)
    _kicker_title(s, "Operating picture", "Three phases. One book.")
    phases = [
        ("Phase 1  ·  Now", "Generate", "Compile source, docs, tickets, and paired logs into draft signatures. Human review stamps truth."),
        ("Phase 2  ·  Later", "107 triage", "Completeness gate first. Missing OFP, SN, or logs become callback questions — not a diagnosis."),
        ("Phase 3  ·  Floor", "FSR log review", "Decode → match signatures → LLM on residue → one kill-check. FSR remains the authority."),
    ]
    for i, (k, t, b) in enumerate(phases):
        x = 0.6 + i * 4.15
        top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(x), PptInches(2.2), PptInches(3.95), PptInches(0.12))
        _set_fill(top, P_STEEL if i == 0 else P_INK)
        sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptInches(x), PptInches(2.32), PptInches(3.95), PptInches(4.2))
        _set_fill(sh, P_ELEV)
        _set_line(sh, P_RULE, 0.75)
        _add_text(s, x + 0.25, 2.55, 3.45, 0.35, k, size=12, bold=True, color=P_STEEL)
        _add_text(s, x + 0.25, 3.0, 3.45, 0.9, t, size=26, bold=True, color=P_INK, font="Georgia")
        _add_text(s, x + 0.25, 4.05, 3.45, 2.1, b, size=15, color=P_MUTED)
    _footer(s, 5, total)

    # 6 Phase 1 flow
    s = _blank(prs)
    _bg(s)
    _kicker_title(
        s,
        "Phase 1  ·  Brain Book generation",
        "The LLM is a compiler, not an oracle.",
        "Source names events. Tickets supply the story. Paired logs create the signature.",
    )
    flow = [
        ("Source + docs", "BIT enums, log strings, ICD names"),
        ("Tickets", "Symptom, config, tries, close"),
        ("Good / bad logs", "Same mode and OFP. Diff is the seed."),
        ("Local compiler", "Decode + diff, then Nemotron drafts YAML"),
        ("Reviewed book", "Maturity: draft → lab / FT reviewed"),
    ]
    for i, (t, b) in enumerate(flow):
        x = 0.45 + i * 2.56
        fill = P_INK if i == 4 else (P_STEEL if i == 3 else P_ELEV)
        tc = P_PAPER if i >= 3 else P_INK
        bc = PptRGB(0xC9, 0xC2, 0xB2) if i >= 3 else P_MUTED
        sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptInches(x), PptInches(2.5), PptInches(2.38), PptInches(3.3))
        _set_fill(sh, fill)
        _set_line(sh, P_RULE if i < 3 else fill, 0.75)
        _add_text(s, x + 0.12, 2.7, 2.14, 0.3, f"0{i+1}", size=11, bold=True, color=P_STEEL if i < 3 else P_PAPER)
        _add_text(s, x + 0.12, 3.1, 2.14, 0.9, t, size=16, bold=True, color=tc, font="Georgia")
        _add_text(s, x + 0.12, 4.15, 2.14, 1.3, b, size=13, color=bc)
    _footer(s, 6, total)

    # 7 Phase 1 capture
    s = _blank(prs)
    _bg(s)
    _kicker_title(s, "Phase 1  ·  Capture rule", "No ticket closed without a log — or a reason.")
    left = [
        ("Do first", "Grep the log/BIT table in code. You cannot name events the software never prints."),
        ("Gold pair", "Ticket + bad log + good log from the same mode and OFP."),
        ("Refuse", "Do not invent limits or root cause from source alone. Open questions stay open."),
    ]
    for i, (t, b) in enumerate(left):
        _card(s, 0.6, 2.15 + i * 1.5, 7.4, 1.38, t, b)
    ink = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptInches(8.25), PptInches(2.15), PptInches(4.5), PptInches(4.38))
    _set_fill(ink, P_INK)
    _add_text(s, 8.5, 2.45, 4.05, 0.4, "Close-the-anomaly minimum", size=14, bold=True, color=P_STEEL)
    _add_text(
        s,
        8.5,
        3.0,
        4.05,
        3.2,
        "1. Config (SN, OFP, ICD)\n2. Symptom in one sentence\n3. Log excerpt — not a 2 GB dump\n4. What failed\n5. What fixed it, or still open\n6. 5–20 lines that would let someone else recognize it",
        size=14,
        color=P_PAPER,
    )
    _footer(s, 7, total)

    # 8 Phase 2
    s = _blank(prs)
    _bg(s)
    _kicker_title(
        s,
        "Phase 2  ·  107 triage",
        "Verify missing data before you diagnose.",
        "A thin field report is not a root-cause problem. It is a completeness problem.",
    )
    steps = [
        ("Inbound 107", "Free text, attachments, sometimes a log."),
        ("Completeness gate", "Extract header. Flag blockers: OFP, SN, last good, already tried, impact."),
        ("Context pack", "Retrieve signatures only for the stated OFP / block / LRU."),
        ("Packet", "Questions if incomplete. Else top three hypotheses with kill-checks."),
    ]
    for i, (t, b) in enumerate(steps):
        y = 2.2 + i * 1.1
        n = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptInches(0.6), PptInches(y), PptInches(0.7), PptInches(0.9))
        _set_fill(n, P_STEEL if i == 1 else P_INK)
        _add_text(s, 0.6, y + 0.22, 0.7, 0.5, f"0{i+1}", size=16, bold=True, color=P_PAPER, align=PP_ALIGN.CENTER)
        body = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptInches(1.5), PptInches(y), PptInches(11.2), PptInches(0.9))
        _set_fill(body, P_ELEV)
        _set_line(body, P_RULE, 0.75)
        _add_text(s, 1.75, y + 0.08, 10.7, 0.35, t, size=16, bold=True, color=P_INK)
        _add_text(s, 1.75, y + 0.44, 10.7, 0.38, b, size=13, color=P_MUTED)
    _footer(s, 8, total)

    # 9 Completeness checklist
    s = _blank(prs)
    _bg(s)
    _kicker_title(s, "Phase 2  ·  Completeness set", "If these are blank, the output is questions.")
    checks = [
        "Unit, site, POC",
        "MDS / NSN / part / SN",
        "SW / OFP / TO in use",
        "First seen vs last known good",
        "Mission impact",
        "Symptom + BIT / fault code",
        "What they already tried",
        "Log, or explicit no-log reason",
    ]
    for i, c in enumerate(checks):
        x = 0.6 + (i % 2) * 6.3
        y = 2.15 + (i // 2) * 1.1
        sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptInches(x), PptInches(y), PptInches(6.05), PptInches(0.95))
        _set_fill(sh, P_ELEV)
        _set_line(sh, P_RULE, 0.75)
        n = s.shapes.add_shape(MSO_SHAPE.OVAL, PptInches(x + 0.2), PptInches(y + 0.26), PptInches(0.42), PptInches(0.42))
        _set_fill(n, P_INK)
        _add_text(s, x + 0.2, y + 0.3, 0.42, 0.35, str(i + 1), size=12, bold=True, color=P_PAPER, align=PP_ALIGN.CENTER)
        _add_text(s, x + 0.8, y + 0.28, 5.0, 0.45, c, size=16, color=P_INK)
    _footer(s, 9, total)

    # 10 Phase 3
    s = _blank(prs)
    _bg(s)
    _kicker_title(
        s,
        "Phase 3  ·  FSR troubleshooting",
        "Parse logs. Match signatures. One next check.",
    )
    flow = [
        ("Dump", "BIT, MC text, syslog attached to the case."),
        ("Decode", "Parser per family → time, source, code, raw."),
        ("Match", "Codes, sequences, OFP filters vs reviewed SIGs."),
        ("Residue", "Nemotron names unmatched lines, fills IS/IS-NOT."),
        ("Action", "One check. Expected if true. Result that kills it."),
    ]
    for i, (t, b) in enumerate(flow):
        x = 0.45 + i * 2.56
        fill = P_INK if i == 4 else (P_STEEL if i == 3 else P_ELEV)
        tc = P_PAPER if i >= 3 else P_INK
        bc = PptRGB(0xC9, 0xC2, 0xB2) if i >= 3 else P_MUTED
        sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptInches(x), PptInches(2.35), PptInches(2.38), PptInches(3.5))
        _set_fill(sh, fill)
        _add_text(s, x + 0.14, 2.55, 2.1, 0.3, f"0{i+1}", size=11, bold=True, color=P_STEEL if i < 3 else P_PAPER)
        _add_text(s, x + 0.14, 3.0, 2.1, 0.7, t, size=18, bold=True, color=tc, font="Georgia")
        _add_text(s, x + 0.14, 3.8, 2.1, 1.7, b, size=13, color=bc)
    _footer(s, 10, total)

    # 11 Constructs
    s = _blank(prs)
    _bg(s)
    _kicker_title(s, "Phase 3  ·  RCA constructs", "Use a short ladder. Do not offer six methodologies.")
    ladder = [
        ("1", "IS / IS-NOT", "What, where, when, extent. Highest value on a thin report."),
        ("2", "Change analysis", "OFP, cable, site, season, TCTO, last known good."),
        ("3", "Signature match", "Events versus baselines. Deterministic."),
        ("4", "Fishbone with status", "Bones: untested / supported / ruled-out. Hypothesis board, not evidence."),
        ("5", "One kill-check", "Then 5 Whys only on the surviving branch."),
    ]
    for i, (n, t, b) in enumerate(ladder):
        y = 2.1 + i * 0.9
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, PptInches(0.7), PptInches(y + 0.08), PptInches(0.5), PptInches(0.5))
        _set_fill(circ, P_INK)
        _add_text(s, 0.7, y + 0.14, 0.5, 0.4, n, size=14, bold=True, color=P_PAPER, align=PP_ALIGN.CENTER)
        _add_text(s, 1.45, y, 4.2, 0.4, t, size=16, bold=True, color=P_INK)
        _add_text(s, 5.8, y, 6.8, 0.7, b, size=14, color=P_MUTED)
    _footer(s, 11, total)

    # 12 Stack
    s = _blank(prs)
    _bg(s)
    _kicker_title(s, "Demo stack", "Open, local, small. Nemotron on the box.")
    rows = [
        ("Model", "NVIDIA Nemotron 3.5 Lightning via Ollama", "Open weights, agentic, ~30B / ~3B active"),
        ("UI", "Open WebUI", "No custom application in v1"),
        ("Decode", "Python parsers + ripgrep", "Logs and source stay deterministic"),
        ("Store", "Git Markdown / YAML", "Reviewable, versioned to OFP"),
        ("Out of scope", "Cloud LLM, Wiki, LGTM, MCP mesh", "Add only after signatures are trusted"),
    ]
    hdr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(0.6), PptInches(2.1), PptInches(12.1), PptInches(0.5))
    _set_fill(hdr, P_INK)
    _add_text(s, 0.8, 2.18, 2.4, 0.35, "LAYER", size=11, bold=True, color=P_PAPER)
    _add_text(s, 3.3, 2.18, 5.2, 0.35, "CHOICE", size=11, bold=True, color=P_PAPER)
    _add_text(s, 8.6, 2.18, 3.8, 0.35, "WHY", size=11, bold=True, color=P_PAPER)
    for i, (a, b, c) in enumerate(rows):
        y = 2.6 + i * 0.78
        bgc = P_ELEV if i % 2 == 0 else P_PAPER
        rct = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(0.6), PptInches(y), PptInches(12.1), PptInches(0.78))
        _set_fill(rct, bgc)
        _add_text(s, 0.8, y + 0.2, 2.4, 0.4, a, size=14, bold=True, color=P_INK)
        _add_text(s, 3.3, y + 0.2, 5.2, 0.4, b, size=14, color=P_INK)
        _add_text(s, 8.6, y + 0.2, 3.8, 0.4, c, size=13, color=P_MUTED)
    _footer(s, 12, total)

    # 13 Ask
    s = _blank(prs)
    _bg(s)
    _kicker_title(s, "Sponsorship ask", "Fund the compiler, not a platform.")
    asks = [
        ("People", "2–3 part-time\n8–12 weeks"),
        ("Hardware", "One existing\nGPU workstation"),
        ("Corpus", "One subsystem, one OFP\n~20 tickets with logs"),
        ("Gate", "10 reviewed signatures\nengineers recognize"),
    ]
    for i, (t, b) in enumerate(asks):
        x = 0.6 + i * 3.15
        sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptInches(x), PptInches(2.2), PptInches(3.0), PptInches(2.6))
        _set_fill(sh, P_INK if i == 3 else P_ELEV)
        tc = P_PAPER if i == 3 else P_INK
        bc = PptRGB(0xC9, 0xC2, 0xB2) if i == 3 else P_MUTED
        _add_text(s, x + 0.2, 2.4, 2.6, 0.35, t.upper(), size=11, bold=True, color=P_STEEL if i != 3 else P_STEEL)
        _add_text(s, x + 0.2, 2.9, 2.6, 1.6, b, size=18, bold=True, color=tc, font="Georgia")
    _add_text(
        s,
        0.7,
        5.15,
        12,
        1.4,
        "Success is not “the model found root cause.” Success is: the eventual cause was already in the top three after decode, or the FSR ran the distinguishing test the same day and the case closed with a new reviewed signature instead of a verbal.",
        size=15,
        color=P_MUTED,
    )
    _footer(s, 13, total)

    # 14 Guardrails
    s = _blank(prs)
    _bg(s, P_INK)
    _add_text(s, 0.7, 0.5, 12, 0.3, "NON-GOALS", size=12, bold=True, color=P_STEEL)
    _add_text(s, 0.7, 0.9, 12, 0.8, "Advise. Do not dispose.", size=36, bold=True, color=P_PAPER, font="Georgia")
    nos = [
        "No cloud tokens. Source, tickets, and logs stay on the box.",
        "No official 107 / JDRS write-back.",
        "No commanding hardware or live aircraft interfaces.",
        "No inventing limits that are not in retrieved text.",
        "The signed procedure and the FSR remain authoritative.",
    ]
    for i, line in enumerate(nos):
        _add_text(s, 0.7, 2.1 + i * 0.7, 12, 0.55, f"  {line}", size=18, color=P_PAPER)
    _add_text(s, 0.7, 6.6, 12, 0.35, "The Brain Book is memory. Disposition stays human.", size=14, color=P_MUTED)

    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))


def main():
    OUT_PUBLIC.mkdir(parents=True, exist_ok=True)
    OUT_ART.mkdir(parents=True, exist_ok=True)
    wp = "Brain-Book-Whitepaper.docx"
    memo = "Brain-Book-Sponsorship-Memo.docx"
    deck = "Brain-Book-Phases.pptx"
    build_whitepaper(OUT_PUBLIC / wp)
    build_memo(OUT_PUBLIC / memo)
    build_deck(OUT_PUBLIC / deck)
    import shutil

    for name in (wp, memo, deck):
        shutil.copy2(OUT_PUBLIC / name, OUT_ART / name)
    print("wrote", wp, memo, deck)


if __name__ == "__main__":
    main()
